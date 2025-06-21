import os
import ast
import re
import pandas as pd
from openpyxl import load_workbook
from openpyxl.styles import Alignment, Font, PatternFill, Border, Side

import copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm, Pt
from PIL import Image

def summary_table(test=False):
    # ===== 這裡設定你的資料夾總路徑 =====
    root_folder = "./record/"  + inference_name + '/'
    data_folder = root_folder + "inference/"
    # 最後要儲存的資料
    records = []

    # 遍歷 sys_1, sys_2, ..., sys_n
    for sys_folder in sorted(os.listdir(data_folder)):
        sys_path = os.path.join(data_folder, sys_folder)

        sys_num = int(re.findall(r'\d+', sys_folder)[0])
        if test and sys_num < 10000:
            continue

        if not os.path.isdir(sys_path):
            continue

        for kpp_folder in os.listdir(sys_path):
            kpp_path = os.path.join(sys_path, kpp_folder)
            if not os.path.isdir(kpp_path):
                continue

            for kvp_folder in os.listdir(kpp_path):
                kvp_path = os.path.join(kpp_path, kvp_folder)
                if not os.path.isdir(kvp_path):
                    continue

                param_file = os.path.join(kvp_path, 'param.txt')
                reward_file = os.path.join(kvp_path, 'reward_info.txt')

                if not (os.path.exists(param_file) and os.path.exists(reward_file)):
                    continue

                try:
                    # 讀 param.txt 的最後一行
                    with open(param_file, 'r') as f:
                        lines = f.readlines()

                        line = lines[0].strip()
                        ini_param = list(map(float, re.findall(r"[-+]?\d*\.\d+|\d+", line)))

                        line = lines[-1].strip()
                        last_param = list(map(float, re.findall(r"[-+]?\d*\.\d+|\d+", line)))

                    # 讀 reward_info.txt 的最後一行
                    with open(reward_file, 'r') as f:
                        lines = f.readlines()

                        ini_reward = ast.literal_eval(lines[0].strip())
                        ini_settling_time = ini_reward['settling_time']
                        ini_overshoot     = ini_reward['overshoot']
                        ini_gm_velocity   = ini_reward['GM_velocity']

                        last_reward = ast.literal_eval(lines[-1].strip())
                        last_settling_time = last_reward['settling_time']
                        last_overshoot     = last_reward['overshoot']
                        last_gm_velocity   = last_reward['GM_velocity']

                    # 提取 sys號碼
                    sys_num = int(re.findall(r'\d+', sys_folder)[0])

                    # 存一筆記錄
                    record = {
                        # 'sys': sys_folder,
                        'sys_num': sys_num,
                        'ini_Kpp': kpp_folder,
                        'ini_Kvp': kvp_folder,
                        'final_kpp': last_param[0],
                        'final_kvp': last_param[1],
                        'ini_settling_time' : ini_settling_time,
                        'last_settling_time' : last_settling_time,
                        'ini_overshoot' : ini_overshoot,
                        'last_overshoot' : last_overshoot,
                        'ini_GM_velocity' : ini_gm_velocity,
                        'last_GM_velocity' : last_gm_velocity
                    }
                    records.append(record)

                except Exception as e:
                    print(f"[錯誤] 資料夾 {kvp_path} 出現問題: {e}")

    # 整理成 DataFrame
    df_main = pd.DataFrame(records)
    df_main['ini_Kpp'] = df_main['ini_Kpp'].apply(lambda x: int(re.findall(r'\d+', str(x))[0]))
    df_main['ini_Kvp'] = df_main['ini_Kvp'].apply(lambda x: int(re.findall(r'\d+', str(x))[0]))

    current_dir = os.path.dirname(__file__)
    csv_path = os.path.join(current_dir, "..", "..", "system_parameters", "system_parameters_table.csv")
    csv_path = os.path.abspath(csv_path)
    df_best = pd.read_csv(csv_path)
    print(csv_path)

    # 只保留需要的欄位
    df_best = df_best[['sys_num', 'best_Kpp', 'best_Kvp', 'settling_time', 'overshoot', 'GM_velocity']].copy()

    # 合併兩個表格
    df_merged = pd.merge(df_main, df_best, on='sys_num', how='left')

    # 重新安排欄位順序（把 best_Kpp, best_Kvp 插到 param_2 後面）
    cols = list(df_merged.columns)
    param2_idx = cols.index('final_kvp')
    new_cols = cols[:param2_idx+1] + ['best_Kpp', 'best_Kvp'] + cols[param2_idx+1:-2] + ['overshoot', 'GM_velocity']
    df_merged = df_merged[new_cols]

    # 依照 sys_num 排序
    df_merged = df_merged.sort_values(by='sys_num')

    # 先讀你的CSV
    if not test:
        output_xlsx = 'result_summary.xlsx'
    else:
        output_xlsx = 'result_summary_test.xlsx'

    df_merged.to_excel(root_folder + output_xlsx, index=False)

    # 開啟 xlsx 檔案
    wb = load_workbook(root_folder + output_xlsx)
    ws = wb.active

    # 置中
    center_alignment = Alignment(horizontal='center', vertical='center')

    # 顏色
    green_font = Font(color="006400")  # 深綠色
    red_font = Font(color="FF0000")    # 紅色

    # 粗線條
    thin_border = Border(bottom=Side(style='thin'))
    gray_fill = PatternFill(start_color="DDDDDD", end_color="EEEEEE", fill_type="solid")  # 淺灰色背景
    white_fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")

    # ======== 處理每一格 ========

    # 先找欄位位置
    header = {cell.value: idx for idx, cell in enumerate(ws[1], start=1)}
    col_ini_settling = header['ini_settling_time']
    col_last_settling = header['last_settling_time']
    col_last_overshoot = header['last_overshoot']
    col_last_gm_velocity = header['last_GM_velocity']
    col_sys_num = header['sys_num']

    # 讀取sys_num做分隔線參考
    previous_sys_num = None
    previous_row = None

    for column in ws.columns:
        max_length = 0
        col_letter = column[0].column_letter  # 取得欄位字母，例如 A, B, C

        for cell in column:
            try:
                # 統一置中
                cell.alignment = center_alignment

                # 找最大字串長度
                if cell.value:
                    cell_length = len(str(cell.value))
                    if cell_length > max_length:
                        max_length = cell_length
            except:
                pass

        # 調整欄寬，給一些額外空間
        adjusted_width = max_length + 2
        ws.column_dimensions[col_letter].width = adjusted_width

    for row in ws.iter_rows(min_row=2):
        sys_num = row[col_sys_num-1].value
        ini_settling = row[col_ini_settling-1].value
        last_settling = row[col_last_settling-1].value
        last_overshoot = row[col_last_overshoot-1].value
        last_gm_velocity = row[col_last_gm_velocity-1].value

        # 每格置中
        for cell in row:
            cell.alignment = center_alignment

        # 1. ini_settling_time > last_settling_time → last_settling_time變深綠色
        if ini_settling is not None and last_settling is not None:
            if ini_settling > last_settling:
                row[col_last_settling-1].font = green_font

        # 2. last_overshoot > 5 → last_overshoot變紅色
        if last_overshoot is not None:
            if last_overshoot > 5.0:
                row[col_last_overshoot-1].font = red_font

        # 3. |last_gm_velocity| > 1 → last_gm_velocity變紅色
        if last_gm_velocity is not None:
            if abs(last_gm_velocity-10.0) > 1.0:
                row[col_last_gm_velocity-1].font = red_font

        # 4. 如果 sys_num 變了 → 畫細底線
        if previous_sys_num is not None and sys_num != previous_sys_num:
            for cell in previous_row:
                cell.border = thin_border

        # 5. 根據 sys_num 是奇數或偶數決定背景色
        if sys_num % 2 == 1:
            fill = gray_fill  # sys_num是奇數→ 灰色
        else:
            fill = white_fill  # sys_num是偶數→ 白色

        for cell in row:
            cell.fill = fill

        previous_sys_num = sys_num
        previous_row = row


    # 儲存
    wb.save(root_folder + output_xlsx)

    print(f"✅ 已經自動調整欄寬、置中，並儲存成 {root_folder + output_xlsx}！")

def create_pptx(test=False):
    # ======== 設定 ========
    ppt_template_path = './utils/雙週會議簡報_模板.pptx'
    root_folder = './record/' + inference_name + '/inference/'
    if not test:
        excel_path = './record/' + inference_name + '/result_summary.xlsx'
        output_ppt_path = './record/' + inference_name + '/調機流程範例.pptx'
    else:
        excel_path = './record/' + inference_name + '/result_summary_test.xlsx'
        output_ppt_path = './record/' + inference_name + '/調機流程範例_test.pptx'

    # 定義圖片對應
    img_files = {
        'Kpp': 'Kpp.jpg',
        'Kvp': 'Kvp.jpg',
        'reward': 'reward.jpg',
        'real_GM_velocity': 'real_GM_velocity.jpg',
        'response_with_label': 'response/response_with_label.png',
    }
    img_keys_order = ['response_with_label', 'Kpp', 'reward', 'real_GM_velocity', 'Kvp']

    # ======== 功能函式 ========

    def copy_slide(prs, slide_idx):
        """複製指定頁面"""
        source = prs.slides[slide_idx]
        slide_layout = source.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)

        for shape in source.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        return new_slide

    def clear_existing_pictures(slide):
        """刪掉這一頁上所有的圖片型別shape，並記錄原本圖片位置大小"""
        picture_shapes = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes.append((shape.left, shape.top, shape.width, shape.height))
        for shape in [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]:
            slide.shapes._spTree.remove(shape.element)
        return picture_shapes

    def insert_picture(slide, img_path, left, top, width, height):
        """插入一張圖，保持比例，居中放置"""
        from PIL import Image
        try:
            with Image.open(img_path) as im:
                img_width, img_height = im.size
        except:
            print(f"❌ 無法讀取圖片 {img_path}")
            return

        img_ratio = img_width / img_height
        box_ratio = width / height

        # 縮放，保持比例
        if img_ratio > box_ratio:
            new_width = width
            new_height = width / img_ratio
        else:
            new_height = height
            new_width = height * img_ratio

        new_left = left + (width - new_width) / 2
        new_top = top + (height - new_height) / 2

        picture = slide.shapes.add_picture(img_path, new_left, new_top, new_width, new_height)
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)

    def update_textboxes(slide, sys_num, ini_Kpp, ini_Kvp, last_Kpp, last_Kvp, best_Kpp, best_Kvp, last_GM_velocity):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            original_text = shape.text

            text = shape.text
            text = text.replace('sys_num', f'sys_{sys_num}')
            text = text.replace('ini_Kpp', f'{ini_Kpp:.1f}')
            text = text.replace('ini_Kvp', f'{ini_Kvp:.1f}')
            text = text.replace('last_Kpp', f'{last_Kpp:.1f}')
            text = text.replace('last_Kvp', f'{last_Kvp:.1f}')
            text = text.replace('best_Kpp', f'{best_Kpp:.1f}')
            text = text.replace('best_Kvp', f'{best_Kvp:.1f}')
            text = text.replace('last_GM_velocity', f'{last_GM_velocity:.1f}')
            shape.text = text

            if text != original_text:
                shape.text = text
                # 判斷是不是標題
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.type == 1:
                    # type==1 是標準Title
                    font_size = Pt(28)  # 標題用28pt
                else:
                    font_size = Pt(18)  # 一般文字用18pt

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Times New Roman"
                        run.font.size = font_size

    def update_table(slide, ini_overshoot, last_overshoot, ini_settling, last_settling, ini_gm, last_gm, best_overshoot, best_settling, best_gm):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                first_cell = table.cell(0, 0).text.strip()
                if "最大過衝量\n(%)" in first_cell:
                    table.cell(1, 0).text = f"{ini_overshoot:.3f} >> {last_overshoot:.3f}"
                    table.cell(1, 1).text = f"{ini_settling:.3f} >> {last_settling:.3f}"
                    table.cell(1, 2).text = f"{ini_gm:.3f} >> {last_gm:.3f}"
                elif "Overshoot" in first_cell:
                    table.cell(1, 0).text = f"{best_overshoot:.3f}"
                    table.cell(1, 1).text = f"{best_settling:.3f}"
                    table.cell(1, 2).text = f"{best_gm:.3f}"
                for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = "Times New Roman"
                                    run.font.size = Pt(20)

    # ======== 主程式開始 ========

    df = pd.read_excel(excel_path)
    prs = Presentation(ppt_template_path)

    # 自動找到「調機流程範例」那頁
    target_idx = None
    for idx, slide in enumerate(prs.slides):
        if slide.shapes.title and "調機流程範例" in slide.shapes.title.text:
            target_idx = idx
            break

    if target_idx is None:
        raise ValueError("❌ 找不到標題包含『調機流程範例』的頁面！")

    # 排序：sys_num從小到大
    sys_folders = []
    for sys_folder in os.listdir(root_folder):
        if not os.path.isdir(os.path.join(root_folder, sys_folder)):
            continue
        if sys_folder.startswith('sys_'):
            num = int(sys_folder.split('_')[1])
            sys_folders.append((num, sys_folder))

    sys_folders.sort()  # 按數字小到大排序

    # 遍歷資料夾
    for _, sys_folder in sys_folders:

        sys_num = int(re.findall(r'\d+', sys_folder)[0])
        if test and sys_num < 10000:
            continue

        sys_path = os.path.join(root_folder, sys_folder)
        if not os.path.isdir(sys_path):
            continue

        for kpp_folder in os.listdir(sys_path):
            kpp_path = os.path.join(sys_path, kpp_folder)
            if not os.path.isdir(kpp_path):
                continue

            for kvp_folder in os.listdir(kpp_path):
                kvp_path = os.path.join(kpp_path, kvp_folder)
                if not os.path.isdir(kvp_path):
                    continue

                print(f"✅ 正在處理 {sys_folder}/{kpp_folder}/{kvp_folder}")

                # 複製範例頁
                slide = copy_slide(prs, target_idx)

                sys_num = int(sys_folder.split('_')[1])
                ini_Kpp = int(kpp_folder.split('_')[1])
                ini_Kvp = int(kvp_folder.split('_')[1])

                matched_row = df[
                    (df['sys_num'] == sys_num) &
                    (df['ini_Kpp'] == ini_Kpp) &
                    (df['ini_Kvp'] == ini_Kvp)
                ]
                if matched_row.empty:
                    print(f"⚠️ 找不到資料 {sys_folder}/{kpp_folder}/{kvp_folder}，跳過")
                    continue
                row = matched_row.iloc[0]

                last_Kpp = float(row['final_kpp'])
                last_Kvp = float(row['final_kvp'])
                best_Kpp = float(row['best_Kpp'])
                best_Kvp = float(row['best_Kvp'])

                ini_overshoot = float(row['ini_overshoot'])
                last_overshoot = float(row['last_overshoot'])
                ini_settling = float(row['ini_settling_time'])
                last_settling = float(row['last_settling_time'])
                ini_gm = float(row['ini_GM_velocity'])
                last_gm = float(row['last_GM_velocity'])

                best_overshoot = float(row['overshoot'])
                best_settling = float(row['settling_time'])
                best_gm = float(row['GM_velocity'])

                update_textboxes(slide, sys_num, ini_Kpp, ini_Kvp, last_Kpp, last_Kvp, best_Kpp, best_Kvp, last_gm)
                update_table(slide, ini_overshoot, last_overshoot, ini_settling, last_settling, ini_gm, last_gm, best_overshoot, best_settling, best_gm)

                # 🔥 只刪除空白或無意義文字框
                for shape in list(slide.shapes):
                    if shape.is_placeholder and shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text in ["", "enter", "Click to add title", "Click to add subtitle"]:  
                            slide.shapes._spTree.remove(shape.element)


                # 修改標題
                for shape in slide.shapes:
                    if shape.has_text_frame and "sys_num" in shape.text:
                        shape.text = f"{sys_folder} / {kpp_folder} / {kvp_folder}"

                # 刪掉舊圖片
                picture_positions = clear_existing_pictures(slide)

                # 插入新圖片（保持比例）
                img_idx = 0
                for key in img_keys_order:
                    filename = img_files[key]
                    img_path = os.path.join(kvp_path, filename)
                    if os.path.exists(img_path) and img_idx < len(picture_positions):
                        left, top, width, height = picture_positions[img_idx]
                        insert_picture(slide, img_path, left, top, width, height)
                        img_idx += 1
                    else:
                        print(f"⚠️ 找不到圖片 {img_path}，跳過")

                # 填表格
                reward_path = os.path.join(kvp_path, 'reward_info.txt')
                if os.path.exists(reward_path):
                    with open(reward_path, 'r', encoding='utf-8') as f:
                        lines = f.readlines()
                        last_data = eval(lines[-1])

                    overshoot = last_data.get('overshoot', 0)
                    settling_time = last_data.get('settling_time', 0)
                    gm_velocity = last_data.get('GM_velocity', 0)

                    for shape in slide.shapes:
                        if shape.has_table:
                            table = shape.table
                            if table.cell(0, 0).text.strip() == "最大過衝量(%)":
                                table.cell(1, 0).text = f"{overshoot:.3f}"
                                table.cell(1, 1).text = f"{settling_time:.3f}"
                                table.cell(1, 2).text = f"{gm_velocity:.3f}"
                            elif table.cell(0, 0).text.strip() == "最大過衝量":
                                table.cell(1, 0).text = f"{overshoot:.3f}"
                                table.cell(1, 1).text = f"{settling_time:.3f}"
                                table.cell(1, 2).text = f"{gm_velocity:.3f}"

    # 儲存簡報
    prs.save(output_ppt_path)
    print(f"✅ 完成！PPT已儲存到 {output_ppt_path}")

if __name__ == "__main__":
    inference_name = 'LSTM_20250608_151201'
    test=False
    # test=True
    summary_table(test=test)
    create_pptx(test=test)